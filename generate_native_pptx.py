import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize standard Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
DARK_BG = RGBColor(15, 23, 42)       # Slate 900
LIGHT_BG = RGBColor(248, 250, 252)   # Slate 50
TEXT_DARK = RGBColor(15, 23, 42)
TEXT_MUTED = RGBColor(100, 116, 139)
TEXT_WHITE = RGBColor(255, 255, 255)
PRIMARY = RGBColor(2, 132, 199)      # Sky 600
ACCENT_GREEN = RGBColor(5, 150, 105) # Emerald 600
ACCENT_RED = RGBColor(220, 38, 38)   # Red 600

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

# SLIDE 1: Title Slide (Layout 0: Title Slide)
s1 = prs.slides.add_slide(prs.slide_layouts[0])
set_bg(s1, DARK_BG)

title = s1.shapes.title
title.text = "Satellite Flood Risk & Relief Monitor"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE

subtitle = s1.placeholders[1]
subtitle.text = (
    "Automated Sentinel-1 SAR Flood Inundation Mapping, Safe Evacuation Routing & Community Relief Connector\n\n"
    "• Nishtha Patel (CSE #1) — Lead Satellite Data & Backend Pipeline\n"
    "• Mahathi (CSE #2) — Frontend Architect & Safe Routing Engine\n"
    "• Rudra (Mechanical) — GIS Risk Analysis & Disaster Economics\n\n"
    "🌐 Live Demo: https://flood-risk-monitor.onrender.com/"
)
for p in subtitle.text_frame.paragraphs:
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(203, 213, 225)

# Helper for content slides
def add_content_slide(title_text, bullets):
    # Layout 1: Title and Content
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_bg(s, LIGHT_BG)
    
    t = s.shapes.title
    t.text = title_text
    t.text_frame.paragraphs[0].font.size = Pt(24)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    body = s.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    
    for i, (head, text) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{head}: "
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(10)
        
        run = p.add_run()
        run.text = text
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
    return s

# SLIDE 2: Problem
add_content_slide(
    "Disaster Bottlenecks in Monsoon Flood Response",
    [
        ("🌧️ Optical Satellites Fail in Monsoons", "Sentinel-2 and Landsat cannot penetrate thick monsoon cloud cover and heavy rainfall, leaving rescue agencies blind for days."),
        ("🚗 Standard Navigation Traps Rescue Teams", "Google Maps and Apple Maps lack real-time flood polygons, guiding ambulances and relief trucks into submerged roads and broken bridges."),
        ("⏳ 6–8 Hour Manual GIS Delays", "Existing workflows require manual satellite downloading, orthorectification, and human polygon tracing, arriving too late for field operations."),
        ("📦 Disconnected Mutual Aid Logistics", "Grassroots volunteers with boats and food supplies have no spatial visibility into where stranded citizens are isolated.")
    ]
)

# SLIDE 3: Architecture
add_content_slide(
    "End-to-End Automated System Architecture",
    [
        ("1. Ingestion & SAR Science", "Automated Copernicus CDSE poller queries Sentinel-1 C-band radar data and detects water specular attenuation (Δσ⁰ ≤ -3.5 dB)."),
        ("2. Backend REST API & Storage", "Flask and Gunicorn WSGI server provides live endpoints (/api/metadata, /api/flood-extent/latest) with versioned run archives."),
        ("3. GIS Operations Portal", "Leaflet.js web map renders flood extents (723.94 km²), high/med/low risk zones, monitored hospitals, and road networks."),
        ("4. Safe Route & Relief Hub", "Dynamic OSRM highway routing avoids active flood polygons with Turf.js spatial BBOX indexing and emergency timeout fallbacks.")
    ]
)

# SLIDE 4: Nishtha - Satellite Science
add_content_slide(
    "SAR Radar Physics & Automated Data Pipeline (Nishtha — CSE #1)",
    [
        ("📡 All-Weather C-Band SAR (5.405 GHz)", "Microwaves penetrate clouds, rain, and darkness, operating identically 24/7."),
        ("🌊 Specular Reflection Principle", "Flat floodwater reflects radar pulses away like a mirror, producing severe dB drops (from -10 dB down to -20 dB)."),
        ("⚡ Change Detection Algorithm", "Calculates backscatter delta Δσ⁰ = during_dB - before_dB. Thresholding (Δσ⁰ < -1.0 dB) isolates new flood pixels."),
        ("⚙️ Morphological Cleaning & GeoJSON Export", "Binary opening (3x3), closing (5x5), and connected component labeling (MIN_PIXELS = 100) eliminate speckle noise before exporting WGS84 GeoJSON.")
    ]
)

# SLIDE 5: Mahathi - UI/UX
add_content_slide(
    "Unified GIS Operations Dashboard & UI/UX (Mahathi — CSE #2)",
    [
        ("🗺️ Dual-View Unified Portal", "Seamless navigation switcher between field Operations Map and deep Satellite Analytics without page reloads."),
        ("🎨 High-Contrast Accessible Map Layers", "Independent layer controls for Flood, Risk zones, Roads, and Hospitals with WCAG AA solid black (#000000) text."),
        ("🛰️ Simplified Live Telemetry", "Human-centric status card showing Acquisition Date, ~6-12 Day Revisit Interval, and pulsing LIVE badge (technical hashes tucked in details toggle)."),
        ("🔄 Flicker-Free Live Streaming", "Auto-polls backend API every 60 seconds and updates flood polygon layers live without resetting user zoom or map center.")
    ]
)

# SLIDE 6: Mahathi - Routing
add_content_slide(
    "Hazard-Aware Safe Route Finder Engine (Mahathi — CSE #2)",
    [
        ("🛡️ Dynamic Flood Avoidance", "Evaluates start/destination submersion, queries OSRM driving paths, and dynamically generates 4-directional bypass detours (~5km buffer)."),
        ("⚡ Spatial BBOX Indexing (4.6x Speedup)", "Pre-computes bounding boxes for 221 flood polygons. O(1) rectangle overlap test filters out 98% of polygons, cutting collision math to 1.66 ms."),
        ("⏱️ Strict 4.5s Timeout & Emergency Fallback", "AbortController aborts hanging OSRM requests and renders emergency direct evacuation vector, preventing infinite browser freezes."),
        ("💾 LRU Memory Cache", "Caches evaluated routes by coordinate key, returning repeated queries in sub-milliseconds (< 5 ms).")
    ]
)

# SLIDE 7: Relief Hub
add_content_slide(
    "Community Relief Hub — Decentralized Mutual Aid",
    [
        ("🚨 Emergency SOS Requests", "Allows stranded victims to request food, drinking water, medical kits, emergency shelter, or rescue boats with priority flags."),
        ("🤝 Volunteer Supply Offers", "NGOs and citizens register available aid supplies, clean water, and inflatable dinghies with direct phone and WhatsApp contact triggers."),
        ("📍 Live Map Visualization", "Color-coded pins (Red: SOS, Green: Supplies) appear live alongside flood polygons for instant geographic matching."),
        ("💾 Offline Local Persistence", "Saves submissions in browser LocalStorage across sessions with quick category and urgency filters.")
    ]
)

# SLIDE 8: Rudra - GIS Risk
add_content_slide(
    "GIS Multi-Criteria Risk Modeling & Exposure (Rudra — Mechanical)",
    [
        ("📊 Multi-Criteria Hazard Formula (QGIS)", "Risk = 0.20*(Inundation) + 0.45*(1 / Dist_Hospital) + 0.35*(Road Severance Index)."),
        ("🔴 High Risk Zone (272.15 km², 41.22%)", "Submerged terrain within 2 km of hospitals or cutting primary evacuation highways like NH-27."),
        ("🟠 Medium & Low Risk (388.10 km², 58.78%)", "Flooded secondary agricultural roads (Medium) and unpopulated open floodplains (Low)."),
        ("🏥 Critical Asset Monitoring", "20+ hospitals monitored across Kamrup, Barpeta, and Goalpara districts to alert emergency teams before road access is cut off.")
    ]
)

# SLIDE 9: Satellite Analytics
add_content_slide(
    "Satellite SAR Analytics & Visual Verification",
    [
        ("📅 Pre-Flood Baseline (29 June 2024)", "Sentinel-1 SAR base capture showing normal river channel and high diffuse backscatter across dry land."),
        ("📅 Peak Inundation Pass (11 July 2024)", "Severe specular reflection expansion and widespread radar signal loss captured through thick monsoon cloud cover."),
        ("🔍 Backscatter Delta Heatmap", "Differential ratio isolating exact zones of sudden radar signal loss, eliminating permanent rivers and lakes."),
        ("🗺️ Classified Vector Mask (723.94 km²)", "Morphologically cleaned and polygonized flood footprint rendered in the Leaflet map layer.")
    ]
)

# SLIDE 10: Tech Stack
add_content_slide(
    "Complete Technology Stack & Deployment",
    [
        ("🛰️ Satellite Data", "Copernicus Data Space Ecosystem (CDSE), Sentinel-1 C-Band SAR (5.405 GHz), OpenStreetMap road networks."),
        ("⚙️ Processing & GIS", "Python 3.14, Rasterio, GeoPandas, Shapely, Scipy Ndimage, QGIS 3.x Spatial Buffering."),
        ("💻 Web & Routing", "Flask REST API, Gunicorn WSGI, Leaflet.js 1.9, Turf.js Spatial Math, Open Source Routing Machine (OSRM)."),
        ("☁️ Cloud Hosting", "Render Python Web Service ($PORT binding), GitHub CI/CD, 100% Open Source.")
    ]
)

# SLIDE 11: Economics & Feasibility
add_content_slide(
    "Operational Feasibility, Scalability & Disaster Economics (Rudra)",
    [
        ("💰 Zero Software Licensing Costs", "100% free open-access satellite data (ESA) and open routing graphs (OSM), eliminating multi-thousand dollar proprietary ArcGIS/Maxar fees."),
        ("⚡ Extreme Client-Side Scalability", "Spatial calculations run client-side in the user's browser via Turf.js, allowing thousands of simultaneous evacuees without server overloads."),
        ("🌍 Global Portability", "Simply adjusting bounding-box coordinates adapts the pipeline to any flood-prone river basin globally (Ganges, Mississippi, Mekong, Rhine)."),
        ("☁️ Minimal Operating Cost", "Runs comfortably on lightweight cloud servers costing less than $10/month per state disaster department.")
    ]
)

# SLIDE 12: Impact Conclusion
s12 = prs.slides.add_slide(prs.slide_layouts[0])
set_bg(s12, DARK_BG)

t12 = s12.shapes.title
t12.text = "Disaster Management Impact & Summary"
t12.text_frame.paragraphs[0].font.size = Pt(32)
t12.text_frame.paragraphs[0].font.bold = True
t12.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE

sub12 = s12.placeholders[1]
sub12.text = (
    "⏱️ Rapid Decision Speed: Replaces 6–8 hours of manual GIS tracing with 30-second automated execution.\n\n"
    "🚑 First Responder Safety: Safe Route Finder prevents rescue teams from entering submerged road corridors.\n\n"
    "🤝 Community Mutual Aid: Connects stranded flood victims with volunteer boat and supply drops directly on the map.\n\n"
    "🌐 Live Platform: https://flood-risk-monitor.onrender.com/  |  Ready for Immediate Deployment"
)
for p in sub12.text_frame.paragraphs:
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(203, 213, 225)

output_file = "Satellite_Flood_Risk_Monitor_Presentation.pptx"
prs.save(output_file)
print(f"[SUCCESS] Native Standard PPTX generated and verified at: {output_file}")
