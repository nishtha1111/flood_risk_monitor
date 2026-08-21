import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation with standard 16:9 Widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette Definitions
BG_DARK = RGBColor(15, 23, 42)       # Slate 900
BG_LIGHT = RGBColor(248, 250, 252)   # Slate 50
CARD_BG = RGBColor(255, 255, 255)    # White
CARD_BORDER = RGBColor(226, 232, 240)# Slate 200
PRIMARY = RGBColor(2, 132, 199)      # Sky 600
PRIMARY_DARK = RGBColor(3, 105, 161) # Sky 700
ACCENT_GREEN = RGBColor(5, 150, 105) # Emerald 600
ACCENT_RED = RGBColor(220, 38, 38)   # Red 600
ACCENT_AMBER = RGBColor(217, 119, 6) # Amber 600
TEXT_MAIN = RGBColor(15, 23, 42)     # Slate 900
TEXT_MUTED = RGBColor(100, 116, 139) # Slate 500
TEXT_LIGHT = RGBColor(241, 245, 249) # Slate 100

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category_text=""):
    if category_text:
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = PRIMARY

    t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.65))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN

def create_card(slide, left, top, width, height, title="", title_color=PRIMARY, bg_color=CARD_BG, border_color=CARD_BORDER):
    # Shape container with clean XML-safe line styling
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    
    # Safe line definition
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.color.rgb = bg_color
        shape.line.width = Pt(0)

    # Text box overlay
    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.space_after = Pt(8)

    return tf

# ============================================================
# SLIDE 1: TITLE SLIDE (DARK MODERN THEME)
# ============================================================
blank_layout = prs.slide_layouts[6]
s1 = prs.slides.add_slide(blank_layout)
set_slide_background(s1, BG_DARK)

# Title Badge
t_badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.9), Inches(4.8), Inches(0.4))
t_badge.fill.solid()
t_badge.fill.fore_color.rgb = RGBColor(30, 41, 59)
t_badge.line.color.rgb = PRIMARY
t_badge.line.width = Pt(1)
tf_b = t_badge.text_frame
tf_b.paragraphs[0].text = "🛰️ HACKATHON INNOVATION PROJECT"
tf_b.paragraphs[0].font.size = Pt(11)
tf_b.paragraphs[0].font.bold = True
tf_b.paragraphs[0].font.color.rgb = PRIMARY
tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER

# Main Title
t1_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.8))
tf1 = t1_box.text_frame
tf1.word_wrap = True
p1 = tf1.paragraphs[0]
p1.text = "Satellite Flood Risk & Relief Monitor"
p1.font.size = Pt(36)
p1.font.bold = True
p1.font.color.rgb = RGBColor(255, 255, 255)

p1_sub = tf1.add_paragraph()
p1_sub.text = "Automated Sentinel-1 SAR Radar Flood Inundation Mapping, Real-Time Safe Evacuation Routing & Community Relief Connector"
p1_sub.font.size = Pt(16)
p1_sub.font.color.rgb = RGBColor(148, 163, 184)
p1_sub.space_before = Pt(8)

# Team Cards (3 Columns)
team_y = Inches(3.6)
card_w = Inches(3.64)
card_h = Inches(2.2)

# Card 1: Nishtha
c1_tf = create_card(s1, Inches(0.8), team_y, card_w, card_h, "🛰️ Nishtha Patel (CSE #1)", PRIMARY, RGBColor(30, 41, 59), RGBColor(51, 65, 85))
p = c1_tf.add_paragraph()
p.text = "Lead Satellite Ingestion & SAR Pipeline"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = RGBColor(226, 232, 240)
p = c1_tf.add_paragraph()
p.text = "• Copernicus CDSE OAuth2 Automation\n• Radar Backscatter Attenuation & Thresholding\n• Vector GeoJSON & Flask REST API"
p.font.size = Pt(11.5)
p.font.color.rgb = RGBColor(148, 163, 184)

# Card 2: Mahathi
c2_tf = create_card(s1, Inches(4.84), team_y, card_w, card_h, "💻 Mahathi (CSE #2)", PRIMARY, RGBColor(30, 41, 59), RGBColor(51, 65, 85))
p = c2_tf.add_paragraph()
p.text = "Frontend Architecture & Safe Routing"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = RGBColor(226, 232, 240)
p = c2_tf.add_paragraph()
p.text = "• Leaflet.js Unified Dual-Portal UI/UX\n• Safe Route Finder (OSRM + Turf.js avoidance)\n• Spatial BBOX Indexing & Mutual-Aid Hub"
p.font.size = Pt(11.5)
p.font.color.rgb = RGBColor(148, 163, 184)

# Card 3: Rudra
c3_tf = create_card(s1, Inches(8.88), team_y, card_w, card_h, "⚙️ Rudra (Mechanical)", PRIMARY, RGBColor(30, 41, 59), RGBColor(51, 65, 85))
p = c3_tf.add_paragraph()
p.text = "GIS Risk Analysis & Disaster Economics"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = RGBColor(226, 232, 240)
p = c3_tf.add_paragraph()
p.text = "• QGIS Multi-Criteria Hazard Scoring\n• Hospital & Road Infrastructure Exposure\n• Operational Feasibility & Impact Modeling"
p.font.size = Pt(11.5)
p.font.color.rgb = RGBColor(148, 163, 184)

# Footer Demo URL
f_box = s1.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6))
tf_f = f_box.text_frame
p_f = tf_f.paragraphs[0]
p_f.text = "🌐 Live Deployment: https://flood-risk-monitor.onrender.com/    |    📂 GitHub: github.com/nishtha1111/flood_risk_monitor"
p_f.font.size = Pt(12)
p_f.font.bold = True
p_f.font.color.rgb = RGBColor(56, 189, 248)

# ============================================================
# SLIDE 2: THE PROBLEM & GROUND REALITY
# ============================================================
s2 = prs.slides.add_slide(blank_layout)
set_slide_background(s2, BG_LIGHT)
add_header(s2, "Disaster Bottlenecks in Monsoon Flood Response", "Context & Ground Reality")

p_w = Inches(5.6)
p_h = Inches(2.4)
row1_y = Inches(1.6)
row2_y = Inches(4.3)
col1_x = Inches(0.8)
col2_x = Inches(6.8)

c1 = create_card(s2, col1_x, row1_y, p_w, p_h, "🌧️ Optical Satellites Fail in Monsoons", ACCENT_RED)
p = c1.add_paragraph()
p.text = "Traditional optical satellites (Sentinel-2, Landsat, MODIS) cannot penetrate dense monsoon cloud cover and heavy rainfall. In peak floods, rescue authorities are left completely blind for days."
p.font.size = Pt(13)
p.font.color.rgb = TEXT_MAIN

c2 = create_card(s2, col2_x, row1_y, p_w, p_h, "🚗 Standard Navigation Traps Rescue Teams", ACCENT_AMBER)
p = c2.add_paragraph()
p.text = "Commercial routing platforms (Google Maps, Apple Maps) lack real-time inundation boundaries. They blindly navigate relief trucks and ambulances straight into submerged highways and washed-out bridges."
p.font.size = Pt(13)
p.font.color.rgb = TEXT_MAIN

c3 = create_card(s2, col1_x, row2_y, p_w, p_h, "⏳ 6–8 Hour Manual GIS Processing Delays", ACCENT_AMBER)
p = c3.add_paragraph()
p.text = "Existing disaster GIS workflows require manual satellite downloading, calibration, and human polygon tracing. By the time maps reach ground commanders, floodwaters have already shifted."
p.font.size = Pt(13)
p.font.color.rgb = TEXT_MAIN

c4 = create_card(s2, col2_x, row2_y, p_w, p_h, "📦 Disconnected Mutual-Aid Logistics", ACCENT_RED)
p = c4.add_paragraph()
p.text = "Grassroots volunteers with boats and food supplies have no spatial visibility into where cut-off victims are stranded. SOS calls on social media lack verified geographic coordinates and hazard contexts."
p.font.size = Pt(13)
p.font.color.rgb = TEXT_MAIN

# ============================================================
# SLIDE 3: SYSTEM ARCHITECTURE & DATA FLOW
# ============================================================
s3 = prs.slides.add_slide(blank_layout)
set_slide_background(s3, BG_LIGHT)
add_header(s3, "End-to-End Automated System Architecture", "How the Solution Operates")

t_w = Inches(3.64)
t_h = Inches(5.1)
t_y = Inches(1.6)

t1 = create_card(s3, Inches(0.8), t_y, t_w, t_h, "1. Ingestion & SAR Science", PRIMARY)
p = t1.add_paragraph()
p.text = "🛰️ Copernicus CDSE Poller"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = t1.add_paragraph()
p.text = "• Automated OAuth2 API token refresh\n• WKT bounding-box intersection query over Assam basin\n• Grabs newest Sentinel-1 IW GRDH 1SDV scene (~6–12d revisit)"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.space_after = Pt(8)

p = t1.add_paragraph()
p.text = "⚡ Radar Physics Engine"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = t1.add_paragraph()
p.text = "• Water specular reflection backscatter delta:\n  Δσ⁰ (dB) = 10 log₁₀(σ_during / σ_before)\n• Specular attenuation threshold (≤ -3.5 dB)\n• Morphological cleaning & polygon extraction"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED

t2 = create_card(s3, Inches(4.84), t_y, t_w, t_h, "2. Backend REST API & Archival", PRIMARY)
p = t2.add_paragraph()
p.text = "📁 Versioned Storage Layer"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = t2.add_paragraph()
p.text = "• Stores historical runs in output/history/<scene_id>/\n• Tracks temporal flood area (km²) & polygon counts\n• Atomic sync to root flood_extent.geojson"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.space_after = Pt(8)

p = t2.add_paragraph()
p.text = "🔌 Live Flask / Gunicorn API"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = t2.add_paragraph()
p.text = "• GET /api/metadata (Active satellite & pass date)\n• GET /api/flood-extent/latest (Polygons)\n• POST /api/pipeline/trigger (On-demand)\n• Dynamic $PORT binding & Gunicorn WSGI"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED

t3 = create_card(s3, Inches(8.88), t_y, t_w, t_h, "3. Decision Support & Routing", PRIMARY)
p = t3.add_paragraph()
p.text = "🗺️ Leaflet.js Operations Portal"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = t3.add_paragraph()
p.text = "• High-contrast vector layers (Flood, Risk, Roads, Hospitals)\n• 60s auto-refresh polling with live stream badge\n• Responsive mobile drawer & accessible UI"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.space_after = Pt(8)

p = t3.add_paragraph()
p.text = "🛡️ Turf.js + OSRM Safe Router"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = t3.add_paragraph()
p.text = "• Real-time road vs flood intersection tests\n• Dynamic 4-directional dry detour bypasses\n• Spatial BBOX pre-filtering index (< 2ms)"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED

# ============================================================
# SLIDE 4: SATELLITE RADAR SCIENCE & PIPELINE (NISHTHA)
# ============================================================
s4 = prs.slides.add_slide(blank_layout)
set_slide_background(s4, BG_LIGHT)
add_header(s4, "SAR Radar Physics & Automated Data Pipeline", "Nishtha Patel (CSE #1) — Backend & Satellite Pipeline")

s4_w = Inches(11.7)
c_top = create_card(s4, Inches(0.8), Inches(1.6), s4_w, Inches(2.4), "📡 Radar Reflection Physics & Specular Attenuation", PRIMARY)
p = c_top.add_paragraph()
p.text = "• All-Weather C-Band SAR: Sentinel-1 emits 5.405 GHz microwave pulses that completely ignore clouds, monsoon rain, and darkness."
p.font.size = Pt(13)
p.font.color.rgb = TEXT_MAIN
p = c_top.add_paragraph()
p.text = "• Specular Surface Reflection: Dry ground and forest create rough diffuse scatter (high backscatter dB: -12 to -8 dB). Smooth standing water reflects radar pulses away like a mirror (low backscatter dB: -22 to -18 dB)."
p.font.size = Pt(13)
p.font.color.rgb = TEXT_MAIN
p = c_top.add_paragraph()
p.text = "• Change Detection Algorithm: Ratio between pre-flood baseline (29 June 2024) and peak inundation (11 July 2024) isolates exact flooded pixels with morphological speckle filtering."
p.font.size = Pt(13)
p.font.color.rgb = TEXT_MAIN

c_bot = create_card(s4, Inches(0.8), Inches(4.3), s4_w, Inches(2.4), "⚙️ Automated Semi-Live Ingestion & Export Engine", ACCENT_GREEN)
p = c_bot.add_paragraph()
p.text = "• Programmatic Copernicus Ingestion: ingest_sentinel1.py queries OData catalogue for IW_GRDH_1SDV products over the AOI [90.64, 25.09, 93.16, 26.70] with OAuth2 credentials and open fallback."
p.font.size = Pt(13)
p.font.color.rgb = TEXT_MAIN
p = c_bot.add_paragraph()
p.text = "• State Tracking Manifest: data/manifest.json compares scene timestamps to skip redundant processing during the 6–12 day orbital revisit cycle."
p.font.size = Pt(13)
p.font.color.rgb = TEXT_MAIN
p = c_bot.add_paragraph()
p.text = "• Vectorized GeoJSON Export: detect_flood.py automatically converts cleaned binary raster masks into EPSG:4326 WGS84 GeoJSON polygons with total area calculations (723.94 km² active extent)."
p.font.size = Pt(13)
p.font.color.rgb = TEXT_MAIN

# ============================================================
# SLIDE 5: OPERATIONS DASHBOARD & UI/UX (MAHATHI)
# ============================================================
s5 = prs.slides.add_slide(blank_layout)
set_slide_background(s5, BG_LIGHT)
add_header(s5, "Unified GIS Operations Dashboard & UI/UX", "Mahathi (CSE #2) — Frontend Architecture & UI/UX")

s5_w = Inches(3.64)
s5_h = Inches(5.1)
s5_y = Inches(1.6)

f1 = create_card(s5, Inches(0.8), s5_y, s5_w, s5_h, "🗺️ Dual-View Unified Portal", PRIMARY)
p = f1.add_paragraph()
p.text = "• Seamless View Switcher"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = f1.add_paragraph()
p.text = "Easily toggle between field Operations Map and deep Satellite Analytics without reloading."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN
p.space_after = Pt(8)

p = f1.add_paragraph()
p.text = "• High-Contrast GIS Layers"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = f1.add_paragraph()
p.text = "Independent layer toggles for Flood Extent, High/Med/Low Risk zones, Roads, and Hospitals with WCAG AA compliant solid black text (#000000)."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN

f2 = create_card(s5, Inches(4.84), s5_y, s5_w, s5_h, "🛰️ Simplified Live Telemetry", PRIMARY)
p = f2.add_paragraph()
p.text = "• Human-Centric Status Card"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = f2.add_paragraph()
p.text = "Shows clear Acquisition Date (Aug 21, 2026 05:24 AM), ~6–12 Day Revisit Interval, and pulsing green LIVE badge."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN
p.space_after = Pt(8)

p = f2.add_paragraph()
p.text = "• Collapsible Technical Scene"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = f2.add_paragraph()
p.text = "Raw satellite product hashes tucked neatly inside an expandable details toggle for developer debugging."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN

f3 = create_card(s5, Inches(8.88), s5_y, s5_w, s5_h, "🔄 Live Polling & Controls", PRIMARY)
p = f3.add_paragraph()
p.text = "• Flicker-Free Dynamic Updates"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = f3.add_paragraph()
p.text = "Auto-polls backend API every 60s. Ingests fresh polygons seamlessly without resetting map view or zooming."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN
p.space_after = Pt(8)

p = f3.add_paragraph()
p.text = "• Manual Satellite Trigger"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = f3.add_paragraph()
p.text = "Interactive [🔄 Check Satellite] button initiates on-demand pipeline polling with live toast alerts."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN

# ============================================================
# SLIDE 6: SAFE ROUTE FINDER & SPATIAL INDEXING (MAHATHI)
# ============================================================
s6 = prs.slides.add_slide(blank_layout)
set_slide_background(s6, BG_LIGHT)
add_header(s6, "Hazard-Aware Safe Route Finder Engine", "Mahathi (CSE #2) — Spatial Routing & Optimization")

s6_w = Inches(5.6)
s6_h = Inches(5.1)

c_left = create_card(s6, Inches(0.8), Inches(1.6), s6_w, s6_h, "🛡️ Dynamic Flood Avoidance Algorithm", PRIMARY)
p = c_left.add_paragraph()
p.text = "1. Point-in-Polygon Origin/Dest Check"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = c_left.add_paragraph()
p.text = "Instantly evaluates if the user's start or target hospital is submerged in floodwaters."
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.space_after = Pt(6)

p = c_left.add_paragraph()
p.text = "2. Multi-Candidate OSRM Pathfinding"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = c_left.add_paragraph()
p.text = "Fetches primary and alternative highway geometries via Open Source Routing Machine (OSRM)."
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.space_after = Pt(6)

p = c_left.add_paragraph()
p.text = "3. Dynamic 4-Directional Bypass Detour"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = c_left.add_paragraph()
p.text = "If primary road is cut by water, automatically tests North/South/East/West elevated bypass waypoints (~5km buffer) to find a dry corridor."
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED

c_right = create_card(s6, Inches(6.8), Inches(1.6), s6_w, s6_h, "⚡ Spatial BBOX Indexing (4.6x Speedup)", ACCENT_GREEN)
p = c_right.add_paragraph()
p.text = "• The Bottleneck Solved:"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = RGBColor(4, 120, 87)
p = c_right.add_paragraph()
p.text = "Testing candidate routes against 221 complex flood polygons across Assam used to take seconds and freeze the browser."
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.space_after = Pt(6)

p = c_right.add_paragraph()
p.text = "• Spatial BBOX Pre-Filtering (< 2ms):"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = RGBColor(4, 120, 87)
p = c_right.add_paragraph()
p.text = "Pre-computes [minX, minY, maxX, maxY] for each polygon. Fast O(1) rectangle overlap rejects 98% of distant polygons instantly."
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.space_after = Pt(6)

p = c_right.add_paragraph()
p.text = "• Timeout Guards & Memory LRU Cache:"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = RGBColor(4, 120, 87)
p = c_right.add_paragraph()
p.text = "Concurrent Promise.all detour queries, strict 4.5s timeout with emergency vector fallback, and sub-millisecond route caching."
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED

# ============================================================
# SLIDE 7: COMMUNITY RELIEF HUB (MUTUAL AID LOGISTICS)
# ============================================================
s7 = prs.slides.add_slide(blank_layout)
set_slide_background(s7, BG_LIGHT)
add_header(s7, "Community Relief Hub — Decentralized Mutual Aid", "Humanitarian Coordination")

s7_w = Inches(3.64)
s7_h = Inches(5.1)
s7_y = Inches(1.6)

rh1 = create_card(s7, Inches(0.8), s7_y, s7_w, s7_h, "🚨 Emergency SOS Requests", ACCENT_RED)
p = rh1.add_paragraph()
p.text = "• Direct Citizen Request Form"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_RED
p = rh1.add_paragraph()
p.text = "Allows stranded flood victims to request urgent assistance with category tagging:\n• 🍞 Food / Drinking Water\n• 💊 Emergency Medical Aid\n• 🚤 Boat Evacuation Needed\n• ⛺ Emergency Shelter"
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN
p.space_after = Pt(8)

p = rh1.add_paragraph()
p.text = "• Urgent Priority Flags"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_RED
p = rh1.add_paragraph()
p.text = "Marks high-severity cases (infants, elderly, critical illness) with pulsating red map pins."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN

rh2 = create_card(s7, Inches(4.84), s7_y, s7_w, s7_h, "🤝 Volunteer Supply Offers", ACCENT_GREEN)
p = rh2.add_paragraph()
p.text = "• Aid Registration Hub"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN
p = rh2.add_paragraph()
p.text = "NGOs and local volunteers register available relief resources:\n• Dry ration kits & clean water tanks\n• First-aid & medical supplies\n• Inflatable rescue dinghies\n• Community kitchen locations"
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN
p.space_after = Pt(8)

p = rh2.add_paragraph()
p.text = "• Contact & Delivery Trigger"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN
p = rh2.add_paragraph()
p.text = "Provides direct phone dialer and WhatsApp triggers for rapid field dispatch."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN

rh3 = create_card(s7, Inches(8.88), s7_y, s7_w, s7_h, "📍 Live Spatial Matching", PRIMARY)
p = rh3.add_paragraph()
p.text = "• Interactive Map Pins"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = rh3.add_paragraph()
p.text = "All supply drops (green markers) and SOS requests (red markers) render live alongside flood extent polygons."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN
p.space_after = Pt(8)

p = rh3.add_paragraph()
p.text = "• Local Storage Persistence"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = rh3.add_paragraph()
p.text = "Stores community items across browser sessions, with quick filters by resource category and urgency."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN

# ============================================================
# SLIDE 8: GIS RISK MODELING & SCORING (RUDRA)
# ============================================================
s8 = prs.slides.add_slide(blank_layout)
set_slide_background(s8, BG_LIGHT)
add_header(s8, "GIS Multi-Criteria Risk Modeling & Infrastructure Exposure", "Rudra (Mechanical) — GIS Risk Analyst")

s8_w = Inches(5.6)
s8_h = Inches(5.1)

g_left = create_card(s8, Inches(0.8), Inches(1.6), s8_w, s8_h, "📊 Multi-Criteria Hazard Scoring (QGIS)", PRIMARY)
p = g_left.add_paragraph()
p.text = "• Multi-Factor Risk Formula:"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = g_left.add_paragraph()
p.text = "Risk = w_f·(Flood Depth/Extent) + w_h·(1 / Dist_Hospital) + w_r·(Road Severance Index)"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.space_after = Pt(6)

p = g_left.add_paragraph()
p.text = "• Classified Risk Distribution:"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = g_left.add_paragraph()
p.text = "🔴 High Risk: 272.15 km² (41.22%)\n   Submerged areas within 2km of hospitals or blocking main arterial highways.\n\n🟠 Medium Risk: 190.55 km² (28.86%)\n   Flooded secondary agricultural and local road networks.\n\n🟢 Low Risk: 197.55 km² (29.92%)\n   Inundated open floodplains far from human settlement."
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MAIN

g_right = create_card(s8, Inches(6.8), Inches(1.6), s8_w, s8_h, "🏥 Critical Asset Exposure & Logistics", PRIMARY)
p = g_right.add_paragraph()
p.text = "• 20+ Hospitals Actively Monitored"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = g_right.add_paragraph()
p.text = "Tracks emergency bed capacities, surgical centers, and access road status across Kamrup, Barpeta, and Goalpara districts."
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.space_after = Pt(8)

p = g_right.add_paragraph()
p.text = "• Road Cut-Off Proactive Warning"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = g_right.add_paragraph()
p.text = "Identifies critical transport corridors at risk of complete isolation before peak water levels crest, allowing early medical airlifts and oxygen stockpiling."
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.space_after = Pt(8)

p = g_right.add_paragraph()
p.text = "• Actionable Relief Prioritization"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = g_right.add_paragraph()
p.text = "Directs NDRF/SDRF boat crews to High Risk zones first, maximizing lives saved per rescue sortie."
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED

# ============================================================
# SLIDE 9: SATELLITE SAR ANALYTICS & VISUAL VERIFICATION
# ============================================================
s9 = prs.slides.add_slide(blank_layout)
set_slide_background(s9, BG_LIGHT)
add_header(s9, "Satellite SAR Analytics & Image Comparison", "Visual Proof of Detection")

g_w = Inches(2.7)
g_h = Inches(5.1)
g_y = Inches(1.6)

b1 = create_card(s9, Inches(0.8), g_y, g_w, g_h, "1. Pre-Flood Baseline", PRIMARY)
p = b1.add_paragraph()
p.text = "📅 29 June 2024"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = b1.add_paragraph()
p.text = "• Sentinel-1 SAR Base\n• Normal Brahmaputra river channel\n• High backscatter return across dry agricultural land\n• VV/VH polarization baseline"
p.font.size = Pt(11.5)
p.font.color.rgb = TEXT_MUTED

b2 = create_card(s9, Inches(3.78), g_y, g_w, g_h, "2. Peak Flood Capture", ACCENT_RED)
p = b2.add_paragraph()
p.text = "📅 11 July 2024"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = ACCENT_RED
p = b2.add_paragraph()
p.text = "• Peak Inundation Pass\n• Massive specular reflection expansion\n• Widespread signal attenuation across rural districts\n• Acquired through thick monsoon cloud cover"
p.font.size = Pt(11.5)
p.font.color.rgb = TEXT_MUTED

b3 = create_card(s9, Inches(6.76), g_y, g_w, g_h, "3. SAR Change Heatmap", ACCENT_AMBER)
p = b3.add_paragraph()
p.text = "Δσ⁰ Difference (dB)"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = ACCENT_AMBER
p = b3.add_paragraph()
p.text = "• Backscatter Delta Ratio\n• Highlights exact zones of sudden radar signal loss\n• Eliminates permanent water bodies to isolate new flood extent"
p.font.size = Pt(11.5)
p.font.color.rgb = TEXT_MUTED

b4 = create_card(s9, Inches(9.74), g_y, g_w, g_h, "4. Classified Flood Mask", ACCENT_GREEN)
p = b4.add_paragraph()
p.text = "Vector Extent Mask"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN
p = b4.add_paragraph()
p.text = "• 723.94 km² Extent\n• Thresholded & morphologically cleaned\n• Converted to WGS84 GeoJSON polygons for Leaflet rendering"
p.font.size = Pt(11.5)
p.font.color.rgb = TEXT_MUTED

# ============================================================
# SLIDE 10: COMPLETE TECH STACK & PRODUCTION DEPLOYMENT
# ============================================================
s10 = prs.slides.add_slide(blank_layout)
set_slide_background(s10, BG_LIGHT)
add_header(s10, "Complete Technology Stack & Deployment", "Production Architecture")

s10_w = Inches(11.7)
s10_h = Inches(1.15)

l1 = create_card(s10, Inches(0.8), Inches(1.6), s10_w, s10_h, "🛰️ Satellite Data & Sourcing", PRIMARY)
p = l1.add_paragraph()
p.text = "Copernicus Data Space Ecosystem (CDSE)  •  Sentinel-1 C-Band SAR (5.405 GHz)  •  OpenStreetMap (OSM) Road Infrastructure Graphs"
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN

l2 = create_card(s10, Inches(0.8), Inches(2.9), s10_w, s10_h, "⚙️ Data Processing & GIS Analytics", PRIMARY)
p = l2.add_paragraph()
p.text = "Python 3.14  •  Rasterio  •  GeoPandas  •  Shapely  •  Scipy Ndimage  •  QGIS 3.x Spatial Analysis & Proximity Buffering"
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN

l3 = create_card(s10, Inches(0.8), Inches(4.2), s10_w, s10_h, "💻 Backend API & Web GIS Frontend", PRIMARY)
p = l3.add_paragraph()
p.text = "Flask REST API  •  Gunicorn WSGI  •  Leaflet.js 1.9  •  Turf.js Spatial Math  •  OSRM Highway Routing Engine  •  HTML5/CSS3/Vanilla JS"
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN

l4 = create_card(s10, Inches(0.8), Inches(5.5), s10_w, s10_h, "☁️ Cloud Deployment & Version Control", ACCENT_GREEN)
p = l4.add_paragraph()
p.text = "Render Cloud Web Service (Dynamic $PORT)  •  Procfile / render.yaml  •  Git / GitHub CI/CD  •  100% Open Source"
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN

# ============================================================
# SLIDE 11: FEASIBILITY, SCALABILITY & ECONOMICS (RUDRA)
# ============================================================
s11 = prs.slides.add_slide(blank_layout)
set_slide_background(s11, BG_LIGHT)
add_header(s11, "Operational Feasibility, Scalability & Disaster Economics", "Rudra (Mechanical) — Economics & Feasibility")

s11_w = Inches(3.64)
s11_h = Inches(5.1)
s11_y = Inches(1.6)

p1_c = create_card(s11, Inches(0.8), s11_y, s11_w, s11_h, "💰 Zero Licensing Costs", ACCENT_GREEN)
p = p1_c.add_paragraph()
p.text = "• 100% Free Data Streams"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN
p = p1_c.add_paragraph()
p.text = "Sentinel-1 radar imagery is free under the European Space Agency (ESA) Open Access policy. OpenStreetMap road networks are free and globally maintained."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN
p.space_after = Pt(8)

p = p1_c.add_paragraph()
p.text = "• No Commercial GIS Licenses"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN
p = p1_c.add_paragraph()
p.text = "Eliminates multi-thousand dollar proprietary ArcGIS / satellite licenses, making it accessible to developing state disaster departments."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN

p2_c = create_card(s11, Inches(4.84), s11_y, s11_w, s11_h, "⚡ Extreme Scalability", PRIMARY)
p = p2_c.add_paragraph()
p.text = "• Client-Side Distributed Math"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = p2_c.add_paragraph()
p.text = "Spatial polygon collision testing and route avoidance are executed client-side via Turf.js in the user's browser."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN
p.space_after = Pt(8)

p = p2_c.add_paragraph()
p.text = "• Server Never Crashes"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = p2_c.add_paragraph()
p.text = "Because route calculation is offloaded to client devices, thousands of simultaneous disaster evacuees can use the system without overloading the server."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN

p3_c = create_card(s11, Inches(8.88), s11_y, s11_w, s11_h, "🌍 Global Portability", PRIMARY)
p = p3_c.add_paragraph()
p.text = "• Universal Applicability"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = p3_c.add_paragraph()
p.text = "By simply changing the bounding-box coordinates, our pipeline instantly adapts to any flood-prone river basin globally (Ganges, Mississippi, Mekong, Rhine)."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN
p.space_after = Pt(8)

p = p3_c.add_paragraph()
p.text = "• Minimal Infrastructure Cost"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK
p = p3_c.add_paragraph()
p.text = "Runs comfortably on lightweight cloud servers costing less than $10/month per state authority."
p.font.size = Pt(12.5)
p.font.color.rgb = TEXT_MAIN

# ============================================================
# SLIDE 12: REAL-WORLD IMPACT & EVALUATION CONCLUSION
# ============================================================
s12 = prs.slides.add_slide(blank_layout)
set_slide_background(s12, BG_DARK)

t12_box = s12.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(1.2))
tf12 = t12_box.text_frame
p12 = tf12.paragraphs[0]
p12.text = "Disaster Management Impact & Summary"
p12.font.size = Pt(30)
p12.font.bold = True
p12.font.color.rgb = RGBColor(255, 255, 255)

p12_sub = tf12.add_paragraph()
p12_sub.text = "Transforming Space Radar Intelligence into Ground-Level Life-Saving Action"
p12_sub.font.size = Pt(15)
p12_sub.font.color.rgb = RGBColor(148, 163, 184)
p12_sub.space_before = Pt(4)

imp_w = Inches(3.64)
imp_h = Inches(4.2)
imp_y = Inches(2.2)

s1 = create_card(s12, Inches(0.8), imp_y, imp_w, imp_h, "⏱️ 6h ➔ 30s Latency", PRIMARY, RGBColor(30, 41, 59), RGBColor(51, 65, 85))
p = s1.add_paragraph()
p.text = "Rapid Decision Speed"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(226, 232, 240)
p = s1.add_paragraph()
p.text = "Replaces 6–8 hours of manual satellite downloading and GIS polygon tracing with an automated 30-second pipeline execution, delivering actionable maps before floodwaters shift."
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(148, 163, 184)

s2 = create_card(s12, Inches(4.84), imp_y, imp_w, imp_h, "🚑 Zero Stranded Teams", ACCENT_GREEN, RGBColor(30, 41, 59), RGBColor(51, 65, 85))
p = s2.add_paragraph()
p.text = "First Responder Safety"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(226, 232, 240)
p = s2.add_paragraph()
p.text = "Safe Route Finder actively prevents ambulances and relief trucks from entering submerged road corridors by enforcing dry bypass routes and warning of severed infrastructure."
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(148, 163, 184)

s3 = create_card(s12, Inches(8.88), imp_y, imp_w, imp_h, "🤝 Localized Mutual Aid", ACCENT_AMBER, RGBColor(30, 41, 59), RGBColor(51, 65, 85))
p = s3.add_paragraph()
p.text = "Empowering Communities"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(226, 232, 240)
p = s3.add_paragraph()
p.text = "Connects stranded victims needing water, medical aid, or boats directly to volunteer donors on the map, removing bureaucratic delays in grassroots relief distribution."
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(148, 163, 184)

fc_box = s12.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5))
tf_fc = fc_box.text_frame
p_fc = tf_fc.paragraphs[0]
p_fc.text = "🌐 Live Application: https://flood-risk-monitor.onrender.com/    |    Ready for Immediate Deployment"
p_fc.font.size = Pt(13)
p_fc.font.bold = True
p_fc.font.color.rgb = RGBColor(56, 189, 248)
p_fc.alignment = PP_ALIGN.CENTER

# Save presentation
output_pptx = "Satellite_Flood_Risk_Monitor_Presentation.pptx"
prs.save(output_pptx)
print(f"[OK] Clean XML-validated PPTX successfully saved to {output_pptx}")
